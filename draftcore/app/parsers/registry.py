from __future__ import annotations

from dataclasses import dataclass
from pathlib import Path

from docx import Document
from pptx import Presentation

from draftcore.app.models.enums import FileType, IngestionStatus


@dataclass(slots=True)
class ParseResult:
    parser_name: str
    status: IngestionStatus
    title: str | None = None
    summary: str | None = None
    searchable_text: str | None = None
    structure_excerpt: str | None = None
    page_count: int | None = None
    paragraph_count: int | None = None


def parse_asset(path: Path, file_type: FileType) -> ParseResult:
    if file_type in {FileType.MD, FileType.TXT}:
        return _parse_text(path)
    if file_type == FileType.DOCX:
        return _parse_docx(path)
    if file_type == FileType.PPTX:
        return _parse_pptx(path)
    if file_type == FileType.IMAGE:
        return _parse_image(path)
    if file_type == FileType.XLSX:
        return ParseResult(
            parser_name="spreadsheet_stub",
            status=IngestionStatus.PENDING,
            title=path.stem,
            summary="Spreadsheet parsing is reserved for a later MVP task.",
        )
    return ParseResult(
        parser_name="unknown_stub",
        status=IngestionStatus.PENDING,
        title=path.stem,
        summary="No parser is available for this file type yet.",
    )


def _parse_text(path: Path) -> ParseResult:
    text = path.read_text(encoding="utf-8", errors="ignore").strip()
    status = IngestionStatus.PARSED if text else IngestionStatus.PARTIAL
    return ParseResult(
        parser_name="text_parser",
        status=status,
        title=_first_line(path.stem, text),
        summary=_summary(text) if text else "Text file is readable but empty.",
        searchable_text=text or None,
        structure_excerpt=_structure_excerpt(text),
        paragraph_count=_paragraph_count(text),
    )


def _parse_docx(path: Path) -> ParseResult:
    document = Document(path)
    paragraphs = [paragraph.text.strip() for paragraph in document.paragraphs if paragraph.text.strip()]
    searchable_text = "\n".join(paragraphs)
    headings = [
        paragraph.text.strip()
        for paragraph in document.paragraphs
        if paragraph.text.strip() and paragraph.style and paragraph.style.name.lower().startswith("heading")
    ]
    status = IngestionStatus.PARSED if searchable_text else IngestionStatus.PARTIAL
    return ParseResult(
        parser_name="docx_parser",
        status=status,
        title=headings[0] if headings else _first_line(path.stem, searchable_text),
        summary=_summary(searchable_text) if searchable_text else "DOCX file has no extractable text.",
        searchable_text=searchable_text or None,
        structure_excerpt="\n".join(headings[:5]) or _structure_excerpt(searchable_text),
        paragraph_count=len(paragraphs) or None,
    )


def _parse_pptx(path: Path) -> ParseResult:
    presentation = Presentation(path)
    slide_titles: list[str] = []
    slide_text_blocks: list[str] = []
    for slide in presentation.slides:
        text_items: list[str] = []
        for shape in slide.shapes:
            if not hasattr(shape, "text"):
                continue
            text = shape.text.strip()
            if text:
                text_items.append(text)
        if text_items:
            slide_titles.append(text_items[0])
            slide_text_blocks.extend(text_items)
    searchable_text = "\n".join(slide_text_blocks)
    status = IngestionStatus.PARSED if searchable_text else IngestionStatus.PARTIAL
    return ParseResult(
        parser_name="pptx_parser",
        status=status,
        title=slide_titles[0] if slide_titles else path.stem,
        summary=_summary(searchable_text) if searchable_text else "PPTX file has no extractable text.",
        searchable_text=searchable_text or None,
        structure_excerpt="\n".join(slide_titles[:5]) or None,
        page_count=len(presentation.slides),
        paragraph_count=len(slide_text_blocks) or None,
    )


def _parse_image(path: Path) -> ParseResult:
    return ParseResult(
        parser_name="image_path_parser",
        status=IngestionStatus.PARSED,
        title=path.stem,
        summary="Image asset is registered by path and filename only in the current MVP.",
        structure_excerpt=path.suffix.lower(),
    )


def _summary(text: str, limit: int = 280) -> str:
    normalized = " ".join(text.split())
    if len(normalized) <= limit:
        return normalized
    return normalized[: limit - 3].rstrip() + "..."


def _structure_excerpt(text: str, limit: int = 5) -> str | None:
    lines = [line.strip() for line in text.splitlines() if line.strip()]
    if not lines:
        return None
    return "\n".join(lines[:limit])


def _paragraph_count(text: str) -> int | None:
    paragraphs = [line for line in text.splitlines() if line.strip()]
    return len(paragraphs) or None


def _first_line(fallback: str, text: str) -> str:
    for line in text.splitlines():
        line = line.strip()
        if line:
            return line
    return fallback
